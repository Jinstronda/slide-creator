"""Core business logic for presentation generation (DRY principle)."""
import os
from pathlib import Path
from typing import Dict, Any
from io import BytesIO

from .excel_parser import get_case_studies
from .ai_selector import select_case_studies, format_selected_for_pptx
from .pptx_generator import (
    generate_presentation,
    add_company_context,
    _replace_in_shapes
)


def generate_presentation_workflow(
    company_name: str,
    company_description: str,
    api_key: str,
    num_cases: int,
    data_path: str,
    template_path: str,
    output_dir: str | None = None
) -> str:
    """
    Generate case study presentation (reusable by CLI/API).

    Returns: Path to generated PPTX file
    """
    _validate_paths(data_path, template_path)

    case_studies = get_case_studies(data_path)
    selected = select_case_studies(
        case_studies,
        company_name,
        company_description,
        api_key,
        num_cases=num_cases
    )

    placeholders = format_selected_for_pptx(selected)
    add_company_context(placeholders, company_name, company_description)

    return generate_presentation(
        template_path,
        placeholders,
        output_dir or "output",
        company_name
    )


def generate_presentation_to_memory(
    company_name: str,
    company_description: str,
    api_key: str,
    num_cases: int,
    data_path: str,
    template_path: str
) -> BytesIO:
    """
    Generate presentation to BytesIO (for API streaming).

    Returns: BytesIO object with PPTX data
    """
    from pptx import Presentation
    from .config import TEMPLATE_CONFIG

    _validate_paths(data_path, template_path)

    case_studies = get_case_studies(data_path)
    selected = select_case_studies(
        case_studies,
        company_name,
        company_description,
        api_key,
        num_cases=num_cases if num_cases > 0 else 4
    )

    placeholders = format_selected_for_pptx(selected)
    add_company_context(placeholders, company_name, company_description)

    prs = Presentation(template_path)
    project_root = Path(__file__).parent.parent

    for slide_idx, slide in enumerate(prs.slides):
        placeholders[TEMPLATE_CONFIG["slide_number"]] = str(slide_idx + 1)
        _replace_in_shapes(slide, placeholders, project_root, slide_idx)

    _filter_slides(prs, num_cases)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _validate_paths(data_path: str, template_path: str) -> None:
    """Validate required files exist."""
    if not os.path.exists(data_path):
        raise FileNotFoundError(f"Data file not found: {data_path}")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")


def _filter_slides(prs, num_cases: int) -> None:
    """
    Remove slides based on presentation type.
    
    Slide mapping:
    - Slide 1 (index 0): 4-case overview
    - Slide 2 (index 1): 2-case grid with CSI
    - Slide 3 (index 2): Single detailed case
    
    num_cases behavior:
    - 0: Keep all slides
    - 1: Keep only Slide 3
    - 2: Keep only Slide 2
    - 4: Keep only Slide 1
    """
    if num_cases == 0:
        return
    
    slides_to_keep = {1: [2], 2: [1], 4: [0]}.get(num_cases, [])
    
    for idx in reversed(range(len(prs.slides._sldIdLst))):
        if idx not in slides_to_keep:
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]
