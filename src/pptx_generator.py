"""Generate PowerPoint presentations from templates."""
import os
import textwrap
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageFont, ImageDraw
from .config import TEMPLATE_CONFIG

# Minimum image size to be considered a case study image (not an icon)
# 1.5 trillion EMUs² ≈ 1.3 x 1.3 inches minimum
MIN_IMAGE_SIZE = 1500000000000

# Standard case study image dimensions (FIXED)
STANDARD_IMAGE_WIDTH_EMU = 2109600  # 5.86 cm = 2.3071 inches
STANDARD_IMAGE_HEIGHT_EMU = 2160000  # 6.0 cm = 2.3622 inches

# Grey box color (RGB) that should hug text
GREY_BOX_RGB = (147, 157, 168)  # 939DA8 in hex

# White box color (RGB) for category labels
WHITE_BOX_RGB = (255, 255, 255)  # FFFFFF in hex

# Text sizing configuration (empirically derived from template measurements)
EMU_PER_POINT = 12700
TITLE_CHAR_WIDTH_FACTOR = 0.56  # Calibrated for ~27 characters per line (TWK Lausanne 11pt)
DESCRIPTION_CHAR_WIDTH_FACTOR = 0.48
TITLE_PADDING_CHARS = 2
MIN_TITLE_WIDTH_PT = 120  # ~1.67 in
MAX_TITLE_LINES = 4
TITLE_HEIGHT_BASE_EMU = 179646  # 0.1965 in
TITLE_HEIGHT_DELTA_EMU = 28444   # 0.0311 in per additional line

# Line positioning (from example1.pptx measurements)
LINE_LEFT_OFFSET_EMU = 85432  # 0.0934in offset from image left
LINE_WIDTH_PERCENT = 0.913  # Line is 91.3% of image width
LINE_TOP_FROM_TITLE_TOP_EMU = 445854  # 0.4876in from title top (FIXED for all titles)
TITLE_FIXED_HEIGHT_EMU = 208000  # Fixed 2-line title height (0.2276in)

# Title vertical alignment (from example1.pptx measurements)
TITLE_FIXED_TOP_EMU = 3638400  # 3.98in from slide top (FIXED for all titles)

# Grey box padding specifications (in EMUs)
# 1 CM = 360000 EMUs (914400 EMUs/inch * 1 inch/2.54 cm)
GREY_BOX_RIGHT_PADDING_EMU = 108000  # 0.3 CM right padding
GREY_BOX_VERTICAL_PADDING_EMU = 46800  # 0.13 CM top/bottom padding

# Character width for grey box text (measured: 0.13 CM per character)
CHAR_WIDTH_EMU = 46800  # 0.13 CM = 46,800 EMUs per character

# Fixed positions for grey boxes on slide 0 (in EMUs from top-left)
# 1 CM = 360,000 EMUs
GREY_BOX_POSITIONS = [
    {'left': int(0.93 * 360000), 'top': int(4.43 * 360000)},  # Case study 1: 0.93 CM, 4.43 CM
    {'left': int(7.07 * 360000), 'top': int(4.42 * 360000)},  # Case study 2: 7.07 CM, 4.42 CM
    {'left': int(13.25 * 360000), 'top': int(4.43 * 360000)},  # Case study 3: 13.25 CM, 4.43 CM
    {'left': int(19.41 * 360000), 'top': int(4.43 * 360000)},  # Case study 4: 19.41 CM, 4.43 CM
]


def generate_presentation(
    template_path: str,
    placeholders: Dict[str, str],
    output_dir: str,
    company_name: str
) -> str:
    """Generate PowerPoint from template with placeholders."""
    prs = Presentation(template_path)
    
    # Get project root directory
    project_root = Path(__file__).parent.parent
    
    for slide_idx, slide in enumerate(prs.slides):
        # Add page number for this slide
        placeholders[TEMPLATE_CONFIG["slide_number"]] = str(slide_idx + 1)
        
        # Use same size for all images (no aspect ratio adjustment)
        _replace_in_shapes(slide, placeholders, project_root, slide_idx)
    
    os.makedirs(output_dir, exist_ok=True)
    safe_name = "".join(c if c.isalnum() else "_" for c in company_name.lower())
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = os.path.join(output_dir, f"{safe_name}_{timestamp}.pptx")
    
    prs.save(output_path)
    return output_path


def _is_grey_box(shape) -> bool:
    """Check if a shape is a grey box that should hug text by checking its name."""
    # Grey boxes are named grey1, grey2, grey3, grey4, grey5, grey6
    if hasattr(shape, 'name') and shape.name:
        return shape.name.lower().startswith('grey')
    return False


def _is_category_box(shape) -> bool:
    """Check if a shape is a white category box that should hug text."""
    if not hasattr(shape, 'fill'):
        return False
    
    try:
        if shape.fill.type == MSO_FILL_TYPE.SOLID:
            rgb = shape.fill.fore_color.rgb
            # Check if color matches white (FFFFFF)
            if (rgb[0], rgb[1], rgb[2]) == WHITE_BOX_RGB:
                # Category boxes are relatively small - filter by size
                # Height should be less than 0.5 inches (457200 EMUs)
                # Width should be less than 3 inches (2742000 EMUs)
                if shape.height < 457200 and shape.width < 2742000:
                    return True
    except:
        pass
    
    return False


def _find_case_study_name_for_grey_box(slide, grey_box) -> Optional:
    """Find the case_study_name shape by matching the grey box number."""
    # Extract the number from grey box name (grey1 -> 1, grey2 -> 2, etc.)
    import re
    match = re.search(r'grey(\d+)', grey_box.name.lower())
    if not match:
        return None
    
    grey_num = match.group(1)
    target_name = f"case_study_{grey_num}_name"
    
    # Find the shape with this name
    for shape in slide.shapes:
        if hasattr(shape, 'name') and shape.name and target_name.lower() in shape.name.lower():
            return shape
    
    return None


def _find_category_text_for_grey_box(slide, grey_box) -> Optional:
    """Find the category text shape that corresponds to this grey box."""
    grey_top = grey_box.top
    grey_left = grey_box.left
    grey_width = grey_box.width
    grey_right = grey_left + grey_width
    
    # Look for category text shapes that overlap horizontally with the grey box
    # and are positioned near it (either on top or very close)
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        if not shape.text:
            continue
        
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape_left + shape.width
        
        # Check horizontal overlap
        horizontal_overlap = not (shape_right < grey_left or shape_left > grey_right)
        
        if horizontal_overlap:
            # Check vertical proximity (shape should be close to grey box)
            vertical_distance = abs(shape_top - grey_top)
            # Only consider shapes within ~0.3 inches (274320 EMUs) vertically
            if vertical_distance < 274320:
                candidates.append((vertical_distance, shape))
    
    # Return the closest shape
    if candidates:
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]
    
    return None


def _find_category_text_on_box(slide, white_box) -> Optional:
    """Find the category text shape that sits on top of the white box."""
    box_top = white_box.top
    box_left = white_box.left
    box_width = white_box.width
    box_height = white_box.height
    box_right = box_left + box_width
    box_bottom = box_top + box_height
    
    # Look for text shapes that overlap with this white box
    # After text replacement, these will be category names like "INFRASTRUCTURE", "HEALTHCARE", etc.
    candidates = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        if not shape.text:
            continue
        
        # Check if shape overlaps with the white box
        shape_left = shape.left
        shape_top = shape.top
        shape_right = shape_left + shape.width
        shape_bottom = shape_top + shape.height
        
        # Check for overlap (must be within or very close to the white box)
        horizontal_overlap = not (shape_right < box_left or shape_left > box_right)
        vertical_overlap = not (shape_bottom < box_top or shape_top > box_bottom)
        
        if horizontal_overlap and vertical_overlap:
            # Calculate how much the text is centered within the box
            overlap_score = abs((shape_left + shape_right)/2 - (box_left + box_right)/2)
            candidates.append((overlap_score, shape))
    
    # Return the most centered text shape
    if candidates:
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]
    
    return None


def _align_title_line_description(slide, slide_idx: int = 0):
    """Align titles, lines (LINE1-4), and descriptions vertically."""
    print(f"\n{'='*80}")
    print(f"DEBUG: ALIGNING TITLES, LINES, DESCRIPTIONS - SLIDE {slide_idx}")
    print(f"{'='*80}")
    
    # Process each case study (1-4)
    for i in range(1, 5):
        title_name = f"case_study_{i}_title"
        line_name = f"line{i}"
        desc_name = f"case_study_{i}_description"
        
        # Find the shapes
        title_shape = None
        line_shape = None
        desc_shape = None
        
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name:
                if title_name.lower() in shape.name.lower():
                    title_shape = shape
                elif line_name.lower() == shape.name.lower():
                    line_shape = shape
                elif desc_name.lower() in shape.name.lower():
                    desc_shape = shape
        
        if not title_shape:
            print(f"\n  No title found for case study {i}")
            continue
        
        print(f"\n--- Aligning Case Study {i} ---")
        
        # Find the corresponding image for this case study by matching order
        # Collect all large images and sort by position (left to right)
        large_images = []
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Check if it's a large image (case study image, not icon)
                if shape.width * shape.height >= MIN_IMAGE_SIZE:
                    large_images.append(shape)
        
        # Sort images by left position (leftmost = first case study)
        large_images.sort(key=lambda s: s.left)
        
        # Match case study i (1-4) with image index i-1 (0-3)
        image_shape = None
        if i <= len(large_images):
            image_shape = large_images[i - 1]
            print(f"  Matched case study {i} with image at position {i-1} (left={image_shape.left/914400:.2f}in)")
        else:
            print(f"  WARNING: No image found for case study {i} (only {len(large_images)} images available)")
        
        # Position title at FIXED vertical position (all titles aligned)
        # Use absolute position from example1.pptx: 3.98in from slide top
        original_title_top = title_shape.top
        new_title_top = TITLE_FIXED_TOP_EMU  # All titles at same vertical position
        
        print(f"  Title: {title_shape.name} at Top={original_title_top/914400:.4f}in -> {new_title_top/914400:.4f}in (FIXED)")
        
        # Set title position BEFORE sizing
        title_shape.top = new_title_top
        
        # Get image width for text block sizing (use fixed width on slide 0)
        image_width = STANDARD_IMAGE_WIDTH_EMU if image_shape else title_shape.width
        
        # Initialize line count and height (fixed for all titles)
        line_count = 2  # Always 2 lines
        calculated_height = TITLE_FIXED_HEIGHT_EMU  # Fixed 2-line height
        
        # Configure title text frame with manual sizing
        if hasattr(title_shape, 'text_frame') and title_shape.text:
            # Get original text first
            title_text = title_shape.text.strip()
            
            # Save original formatting before any changes
            original_formatting = {}
            if title_shape.text_frame.paragraphs:
                para = title_shape.text_frame.paragraphs[0]
                if para.runs:
                    orig_run = para.runs[0]
                    if orig_run.font.size:
                        original_formatting['size'] = orig_run.font.size
                    if orig_run.font.bold is not None:
                        original_formatting['bold'] = orig_run.font.bold
                    if orig_run.font.italic is not None:
                        original_formatting['italic'] = orig_run.font.italic
                    if orig_run.font.color.rgb:
                        original_formatting['color'] = orig_run.font.color.rgb

            # Get font size for wrapping calculation
            font_size_pt = 11  # Default
            if 'size' in original_formatting and original_formatting['size']:
                font_size_pt = original_formatting['size'].pt
            
            # Force all titles to be 2 lines maximum
            wrapped_lines, truncated = _wrap_text_to_lines(
                title_text,
                image_width,
                font_size_pt,
                2,  # ALWAYS 2 lines
                TITLE_CHAR_WIDTH_FACTOR
            )
            final_text = "\n".join(wrapped_lines)
            line_count = 2  # Always 2 lines
            
            if truncated:
                print(f"    WARNING: Title truncated to 2 lines")

            # Configure text frame with manual sizing (NO AUTO-SIZE)
            text_frame = title_shape.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-size completely
            text_frame.vertical_anchor = MSO_ANCHOR.TOP
            text_frame.margin_left = Pt(0)
            text_frame.margin_right = Pt(0)
            text_frame.margin_top = Pt(0)
            text_frame.margin_bottom = Pt(0)
            
            # Set width and align with image horizontally
            title_shape.width = image_width
            if image_shape:
                title_shape.left = image_shape.left
                print(f"    Position: left={title_shape.left/914400:.4f}in (aligned with image)")
            
            # Use fixed 2-line title height for all titles
            calculated_height = TITLE_FIXED_HEIGHT_EMU
            title_shape.height = calculated_height
            
            # Set the text with original formatting
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                para.clear()
                run = para.add_run()
                run.text = final_text
                # Restore original formatting
                if 'size' in original_formatting:
                    run.font.size = original_formatting['size']
                if 'bold' in original_formatting:
                    run.font.bold = original_formatting['bold']
                if 'italic' in original_formatting:
                    run.font.italic = original_formatting['italic']
                if 'color' in original_formatting:
                    run.font.color.rgb = original_formatting['color']
            
            # Title top position already set above (TITLE_FIXED_TOP_EMU)
            # No need to restore, keep the fixed position
            
            print(f"  Title: '{final_text[:50]}...'")
            print(f"    Manual sizing: width={title_shape.width/914400:.4f}in, height={title_shape.height/914400:.4f}in, lines={line_count}")
        
        # Position LINE at fixed distance from title top (matching example1.pptx)
        # Title top to Line top: 0.4876in (445854 EMUs) - FIXED for all titles
        if line_shape:
            # Fixed position from title top
            new_line_top = title_shape.top + LINE_TOP_FROM_TITLE_TOP_EMU
            
            old_line_top = line_shape.top / 914400
            new_line_top_inches = new_line_top / 914400
            
            # Calculate line width and position to match example1.pptx
            new_line_width = _calculate_line_width(image_width)
            old_line_width = line_shape.width / 914400
            new_line_width_inches = new_line_width / 914400
            
            print(f"  LINE: {line_shape.name}")
            print(f"    Old Top: {old_line_top:.2f}in, Width: {old_line_width:.2f}in")
            print(f"    NEW Top: {new_line_top_inches:.2f}in ({LINE_TOP_FROM_TITLE_TOP_EMU/914400:.4f}in from title top - FIXED)")
            print(f"    NEW Width: {new_line_width_inches:.2f}in ({LINE_WIDTH_PERCENT*100:.1f}% of image)")
            
            line_shape.top = new_line_top
            line_shape.width = new_line_width
            
            # Position line with offset from image left (matching example1.pptx)
            if image_shape:
                line_shape.left = image_shape.left + LINE_LEFT_OFFSET_EMU
                print(f"    NEW Left: {line_shape.left/914400:.4f}in (image + {LINE_LEFT_OFFSET_EMU/914400:.4f}in offset)")
        else:
            print(f"  WARNING: No LINE{i} found")
        
        # Position description right below LINE (EXACT spacing from example1.pptx)
        # Line bottom to Desc top: 0.1156in (105713 EMUs)
        if desc_shape and line_shape:
            line_bottom = line_shape.top + line_shape.height
            new_desc_top = line_bottom + 105713  # Exact spacing from example1
            
            old_desc_top = desc_shape.top / 914400
            new_desc_top_inches = new_desc_top / 914400
            
            print(f"  Description: {desc_shape.name}")
            print(f"    Old Top: {old_desc_top:.2f}in")
            print(f"    NEW Top: {new_desc_top_inches:.2f}in (0.1156in below LINE)")
            
            desc_shape.top = new_desc_top
            
            # Align description width with line width (NO AUTO-SIZE)
            if hasattr(desc_shape, 'text_frame'):
                desc_text_frame = desc_shape.text_frame
                desc_text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-size
                desc_text_frame.word_wrap = True
                desc_text_frame.vertical_anchor = MSO_ANCHOR.TOP
                desc_text_frame.margin_left = Pt(0)
                desc_text_frame.margin_right = Pt(0)
                desc_text_frame.margin_top = Pt(0)
                desc_text_frame.margin_bottom = Pt(0)
                # Use same width as line for consistency
                desc_shape.width = new_line_width
                if image_shape:
                    desc_shape.left = image_shape.left + LINE_LEFT_OFFSET_EMU
                
                # Adjust description font size based on text length
                if desc_shape.text and desc_text_frame.paragraphs:
                    desc_text_length = len(desc_shape.text)
                    # Get current font size
                    current_font_size = None
                    for para in desc_text_frame.paragraphs:
                        if para.runs:
                            for run in para.runs:
                                if run.font.size:
                                    current_font_size = run.font.size
                                    break
                            if current_font_size:
                                break
                    
                    if current_font_size:
                        original_pt = current_font_size.pt
                        new_size_pt = original_pt
                        
                        # Adjust font based on description length
                        if desc_text_length < 100:
                            # Very short: increase by 2pt
                            new_size_pt = original_pt + 2
                            print(f"    Description short ({desc_text_length} chars): increased font {original_pt}pt -> {new_size_pt}pt")
                        elif desc_text_length < 150:
                            # Short: increase by 1pt
                            new_size_pt = original_pt + 1
                            print(f"    Description short ({desc_text_length} chars): increased font {original_pt}pt -> {new_size_pt}pt")
                        elif desc_text_length > 200:
                            # Long: reduce to 6pt
                            new_size_pt = 6
                            print(f"    Description long ({desc_text_length} chars): reduced font {original_pt}pt -> {new_size_pt}pt")
                        
                        # Apply new size if changed
                        if new_size_pt != original_pt:
                            for para in desc_text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.size = Pt(new_size_pt)
        elif desc_shape:
            print(f"  WARNING: Found description but no LINE for case study {i}")
            if hasattr(desc_shape, 'text_frame'):
                desc_text_frame = desc_shape.text_frame
                desc_text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Disable auto-size
                desc_text_frame.word_wrap = True
                desc_text_frame.vertical_anchor = MSO_ANCHOR.TOP
                desc_text_frame.margin_left = Pt(0)
                desc_text_frame.margin_right = Pt(0)
                desc_text_frame.margin_top = Pt(0)
                desc_text_frame.margin_bottom = Pt(0)
                # Use calculated width matching line width even without line
                calculated_width = _calculate_line_width(image_width)
                desc_shape.width = calculated_width
                if image_shape:
                    desc_shape.left = image_shape.left + LINE_LEFT_OFFSET_EMU
                
                # Adjust description font size based on text length
                if desc_shape.text and desc_text_frame.paragraphs:
                    desc_text_length = len(desc_shape.text)
                    # Get current font size
                    current_font_size = None
                    for para in desc_text_frame.paragraphs:
                        if para.runs:
                            for run in para.runs:
                                if run.font.size:
                                    current_font_size = run.font.size
                                    break
                            if current_font_size:
                                break
                    
                    if current_font_size:
                        original_pt = current_font_size.pt
                        new_size_pt = original_pt
                        
                        # Adjust font based on description length
                        if desc_text_length < 100:
                            # Very short: increase by 2pt
                            new_size_pt = original_pt + 2
                            print(f"    Description short ({desc_text_length} chars): increased font {original_pt}pt -> {new_size_pt}pt")
                        elif desc_text_length < 150:
                            # Short: increase by 1pt
                            new_size_pt = original_pt + 1
                            print(f"    Description short ({desc_text_length} chars): increased font {original_pt}pt -> {new_size_pt}pt")
                        elif desc_text_length > 200:
                            # Long: reduce to 6pt
                            new_size_pt = 6
                            print(f"    Description long ({desc_text_length} chars): reduced font {original_pt}pt -> {new_size_pt}pt")
                        
                        # Apply new size if changed
                        if new_size_pt != original_pt:
                            for para in desc_text_frame.paragraphs:
                                for run in para.runs:
                                    run.font.size = Pt(new_size_pt)
    
    print(f"{'='*80}\n")


def _align_metrics_with_labels(slide, slide_idx: int = 0):
    """Combine metric numbers with their labels into a single text box."""
    print(f"\n{'='*80}")
    print(f"DEBUG: COMBINING METRICS WITH LABELS - SLIDE {slide_idx}")
    print(f"{'='*80}")
    
    # Find all metric shapes (n1, n2, n3, n4, n5, n6)
    metric_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'name') and shape.name:
            if shape.name.lower() in ['n1', 'n2', 'n3', 'n4', 'n5', 'n6']:
                metric_shapes.append(shape)
    
    print(f"Found {len(metric_shapes)} metric shapes on slide {slide_idx}")
    for ms in metric_shapes:
        print(f"  - {ms.name}")
    
    shapes_to_delete = []
    
    for metric_shape in metric_shapes:
        # First, remove trailing spaces from the metric text
        metric_text_cleaned = None
        if hasattr(metric_shape, 'text_frame') and metric_shape.text:
            for paragraph in metric_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        old_metric_text = run.text
                        run.text = run.text.rstrip()
                        metric_text_cleaned = run.text
                        if old_metric_text != run.text:
                            print(f"\n  Cleaned metric {metric_shape.name}: '{old_metric_text}' -> '{run.text}'")
        
        # Configure text frame for tight fit (will resize later with combined text)
        if metric_text_cleaned and hasattr(metric_shape, 'text_frame'):
            # Disable auto-sizing and set margins to 0 for tight fit
            metric_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE  # Fixed size, no auto-resize!
            metric_shape.text_frame.word_wrap = False  # No word wrap
            metric_shape.text_frame.margin_left = 0
            metric_shape.text_frame.margin_right = 0
            metric_shape.text_frame.margin_top = 0
            metric_shape.text_frame.margin_bottom = 0
        
        # Extract the number (n1 -> 1, n2 -> 2, etc.)
        metric_num = metric_shape.name.lower().replace('n', '')
        label_name = f"metric_label_case_study_{metric_num}"
        
        # Find the label shape
        label_shape = None
        label_text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'name') and shape.name and label_name.lower() in shape.name.lower():
                if hasattr(shape, 'text') and shape.text:
                    label_shape = shape
                    # Get the label text and clean it
                    label_text = shape.text.strip()
                    print(f"\n  Found label for {metric_shape.name}: '{label_text[:40]}'")
                    break
        
        if not label_shape or not label_text:
            print(f"  WARNING: No label text found for {metric_shape.name}")
            continue
        
        print(f"\n--- Combining {metric_shape.name} + label ---")
        print(f"  Metric: '{metric_text_cleaned}'")
        print(f"  Label: '{label_text}'")
        
        # Combine metric and label into metric text box with different font sizes
        if hasattr(metric_shape, 'text_frame') and metric_text_cleaned:
            # Get the metric's font size
            metric_font_size = 30  # Default
            if metric_shape.text_frame.paragraphs:
                for para in metric_shape.text_frame.paragraphs:
                    if para.runs:
                        if para.runs[0].font.size:
                            metric_font_size = para.runs[0].font.size.pt
                        break
            
            # Determine spacing based on metric type
            metric_text = metric_text_cleaned.strip()
            
            if metric_text.startswith('x') or metric_text.startswith('X'):
                # x2.5 type - minimal space
                spaces = "  "
                spacing_desc = "2 spaces (x-type)"
            elif '%' in metric_text:
                # Percentage type - just 1 space
                spaces = " "
                spacing_desc = "1 space (%-type)"
            else:
                # Number with + - just 1 space
                spaces = " "
                spacing_desc = "1 space (number+type)"
            
            print(f"  Creating formatted text: Metric '{metric_text}' @ {metric_font_size}pt + {spacing_desc} + Label '{label_text}' @ 11pt")
            
            # Calculate approximate width of combined text to check for overflow
            # Metric width: metric_font_size * 0.6 * char_count
            metric_width_pt = metric_font_size * 0.6 * len(metric_text)
            # Label width: 11pt * 0.5 * char_count (italic is narrower)
            label_font_size = 11
            label_width_pt = label_font_size * 0.5 * len(label_text)
            # Spaces width
            spaces_width_pt = label_font_size * 0.5 * len(spaces)
            
            # Total width in points
            total_width_pt = metric_width_pt + spaces_width_pt + label_width_pt
            total_width_emu = int(total_width_pt * 12700)
            
            # Check overflow against IMAGE width (for slide 0) or original available width
            if slide_idx == 0:
                # Use fixed image width for slide 0
                max_allowed_width = STANDARD_IMAGE_WIDTH_EMU
            else:
                # For other slides, use metric box + label box combined width
                max_allowed_width = metric_shape.width + label_shape.width
            
            max_allowed_width_pt = max_allowed_width / 12700
            
            print(f"  Calculated text width: {total_width_pt:.1f}pt ({total_width_emu/914400:.2f}in)")
            print(f"  Max allowed width: {max_allowed_width_pt:.1f}pt ({max_allowed_width/914400:.2f}in)")
            
            # If overflow, reduce label font size
            if total_width_emu > max_allowed_width:
                reduction_factor = max_allowed_width / total_width_emu
                label_font_size = int(11 * reduction_factor)
                # Ensure minimum font size of 8pt
                label_font_size = max(8, label_font_size)
                print(f"  OVERFLOW DETECTED! Reducing label font: 11pt -> {label_font_size}pt")
            
            # Clear existing runs and create new ones with different formatting
            paragraph = metric_shape.text_frame.paragraphs[0]
            paragraph.clear()
            
            from pptx.dml.color import RGBColor
            
            # Add metric number (white, BOLD AND ITALIC)
            run1 = paragraph.add_run()
            run1.text = metric_text
            run1.font.size = Pt(metric_font_size)
            run1.font.bold = True
            run1.font.italic = True
            run1.font.color.rgb = RGBColor(255, 255, 255)  # White
            
            # Add spaces (white, italic only)
            run2 = paragraph.add_run()
            run2.text = spaces
            run2.font.size = Pt(label_font_size)
            run2.font.bold = False
            run2.font.italic = True
            run2.font.color.rgb = RGBColor(255, 255, 255)  # White
            
            # Add label (white, ITALIC ONLY - not bold)
            run3 = paragraph.add_run()
            run3.text = label_text
            run3.font.size = Pt(label_font_size)
            run3.font.bold = False  # Label is NOT bold
            run3.font.italic = True
            run3.font.color.rgb = RGBColor(255, 255, 255)  # White
            
            print(f"  Created 3 runs: '{metric_text}' ({metric_font_size}pt white/BOLD/italic) + '{spaces}' + '{label_text}' ({label_font_size}pt white/italic)")
            
            # Expand the metric box to fit the combined text
            # Get original label width to add to metric box
            original_label_width = label_shape.width
            new_combined_width = metric_shape.width + original_label_width
            
            # Ensure metric box doesn't exceed image width on slide 0
            if slide_idx == 0:
                new_combined_width = min(new_combined_width, STANDARD_IMAGE_WIDTH_EMU)
            
            old_width = metric_shape.width / 914400
            new_width = new_combined_width / 914400
            
            print(f"  Expanding metric box: {old_width:.2f}in -> {new_width:.2f}in (added label width)")
            metric_shape.width = new_combined_width
        
        # Mark label shape for deletion
        shapes_to_delete.append(label_shape)
        print(f"  Marked {label_shape.name} for deletion")
    
    # Delete the label shapes
    for shape in shapes_to_delete:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
            print(f"  Deleted shape: {shape.name}")
        except Exception as e:
            print(f"  Warning: Could not delete shape: {e}")
    
    print(f"{'='*80}\n")


def _resize_grey_boxes(slide, slide_idx: int = 0):
    """Resize all grey boxes (grey1, grey2, etc.) to hug their category text."""
    grey_boxes = [shape for shape in slide.shapes if _is_grey_box(shape)]
    
    print(f"\n{'='*80}")
    print(f"DEBUG: GREY BOX RESIZING - SLIDE {slide_idx}")
    print(f"{'='*80}")
    print(f"Found {len(grey_boxes)} grey boxes on slide {slide_idx}")
    
    for grey_box in grey_boxes:
        print(f"\n--- Processing {grey_box.name} ---")
        print(f"  Position: Left={grey_box.left/914400:.2f}in, Top={grey_box.top/914400:.2f}in")
        print(f"  Current size: Width={grey_box.width/914400:.2f}in, Height={grey_box.height/914400:.2f}in")
        
        # Find the case_study_name shape (what should be below the grey box)
        name_shape = _find_case_study_name_for_grey_box(slide, grey_box)
        
        if not name_shape:
            print(f"  WARNING: No case_study_name found!")
            continue
        
        print(f"  Found case_study_name: '{name_shape.text[:50] if name_shape.text else '[empty]'}'")
        print(f"  Name position: Left={name_shape.left/914400:.2f}in, Top={name_shape.top/914400:.2f}in")
        print(f"  Name size: Width={name_shape.width/914400:.2f}in")
        
        # Extract grey box number (grey1 -> 0, grey2 -> 1, etc.)
        import re
        match = re.search(r'grey(\d+)', grey_box.name.lower())
        if not match:
            print(f"  WARNING: Cannot extract number from grey box name")
            continue

        grey_num = int(match.group(1))
        grey_index = grey_num - 1  # Convert to 0-based index

        # Find the category text shape for this grey box
        text_shape = _find_category_text_for_grey_box(slide, grey_box)

        if text_shape and text_shape.text:
            text_content = text_shape.text.strip()
            print(f"  Found category text: '{text_content}'")
            print(f"  Text length: {len(text_content)} characters")

            # Calculate exact text width from character count (0.13 CM per character)
            calculated_text_width = len(text_content) * CHAR_WIDTH_EMU
            print(f"  Calculated text width: {calculated_text_width/914400:.4f}in ({len(text_content)} × 0.13 CM)")

            # Configure text frame
            if hasattr(text_shape, 'text_frame'):
                text_frame = text_shape.text_frame
                # Set margins to zero for tight fit
                text_frame.margin_left = 0
                text_frame.margin_right = 0
                text_frame.margin_top = 0
                text_frame.margin_bottom = 0
                # Disable auto-size - we set dimensions manually
                text_frame.auto_size = MSO_AUTO_SIZE.NONE
                text_frame.word_wrap = False

                print(f"  Set text frame margins to zero, auto_size to NONE")

            # Store original position
            text_left = text_shape.left
            text_top = text_shape.top
            text_height = text_shape.height

            # Set text shape width to calculated width
            text_shape.width = int(calculated_text_width)

            print(f"  Set text shape width: {text_shape.width/914400:.4f}in (calculated from text)")
            print(f"  Text position (keeping): left={text_left/914400:.4f}in, top={text_top/914400:.4f}in")

            # Define padding
            left_padding_emu = 50000  # Minimal left padding (~0.055 inches)
            right_padding_emu = GREY_BOX_RIGHT_PADDING_EMU  # 0.3 CM right padding

            print(f"  Padding: left={left_padding_emu/914400:.4f}in, right={right_padding_emu/914400:.4f}in (0.3 CM)")
            print(f"  Padding: vertical={GREY_BOX_VERTICAL_PADDING_EMU/914400:.4f}in (0.13 CM)")

            # Calculate grey box dimensions from manually sized text
            new_grey_width = calculated_text_width + left_padding_emu + right_padding_emu
            new_grey_height = text_height + (2 * GREY_BOX_VERTICAL_PADDING_EMU)

            # Center grey box around text (equal padding on all sides)
            total_horizontal_padding = left_padding_emu + right_padding_emu
            new_grey_left = text_left - (total_horizontal_padding // 2)

            total_vertical_padding = 2 * GREY_BOX_VERTICAL_PADDING_EMU
            new_grey_top = text_top - (total_vertical_padding // 2)

            print(f"  Centering grey box:")
            print(f"    Horizontal: total={total_horizontal_padding/914400:.4f}in, half={total_horizontal_padding//2/914400:.4f}in")
            print(f"    Vertical: total={total_vertical_padding/914400:.4f}in, half={total_vertical_padding//2/914400:.4f}in")

            # Update grey box
            grey_box.width = int(new_grey_width)
            grey_box.height = int(new_grey_height)
            grey_box.left = int(new_grey_left)
            grey_box.top = int(new_grey_top)

            print(f"  FINAL:")
            print(f"    Grey box: left={grey_box.left/914400:.4f}in, top={grey_box.top/914400:.4f}in")
            print(f"    Grey box: width={grey_box.width/914400:.4f}in, height={grey_box.height/914400:.4f}in")
            print(f"    Text: left={text_shape.left/914400:.4f}in, width={text_shape.width/914400:.4f}in (MANUALLY SIZED)")
        else:
            print(f"  WARNING: No category text found or text is empty!")
    
    print(f"{'='*80}\n")


def _resize_category_boxes(slide, slide_idx: int = 0):
    """Resize white category boxes to hug the category text."""
    category_boxes = [shape for shape in slide.shapes if _is_category_box(shape)]
    
    print(f"\nFound {len(category_boxes)} category boxes on slide {slide_idx}")
    
    for cat_box in category_boxes:
        # Find the category text shape on top of this white box
        text_shape = _find_category_text_on_box(slide, cat_box)
        
        if text_shape and text_shape.text:
            # Calculate text width based on font size
            font_size = 11  # Default
            if text_shape.text_frame.paragraphs:
                for para in text_shape.text_frame.paragraphs:
                    if para.runs:
                        if para.runs[0].font.size:
                            font_size = para.runs[0].font.size.pt
                        break
            
            # Approximate character width: 0.52 * font size in points
            # (based on analysis of example1.pptx category boxes)
            char_width_pt = font_size * 0.52
            char_width_emu = int(char_width_pt * 12700)
            
            # Calculate new width based on actual text content
            text_content = text_shape.text.strip()
            new_width = len(text_content) * char_width_emu
            
            # Add some padding (20 EMUs per side = ~40 total)
            padding = 40 * 12700
            new_width += padding
            
            # Update category box width
            old_width_inches = cat_box.width / 914400
            new_width_inches = new_width / 914400
            
            # Clean text for printing
            clean_text = text_content.encode('ascii', 'ignore').decode('ascii')[:20]
            print(f"  Resizing category box: '{clean_text}' "
                  f"({old_width_inches:.2f}in -> {new_width_inches:.2f}in)")
            
            cat_box.width = new_width


def _replace_in_shapes(slide, placeholders: Dict[str, str], project_root: Path, slide_idx: int = 0):
    """Replace placeholders in all shapes and insert images."""
    shapes_to_remove = []
    images_to_add = []
    
    # Collect all picture shapes on the slide
    picture_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # PICTURE
            picture_shapes.append(shape)
    
    # Filter to only LARGE images (case study placeholders), exclude small icons
    if picture_shapes:
        picture_shapes_with_size = [(shape, shape.width * shape.height) for shape in picture_shapes]
        
        print(f"\n=== ALL IMAGES ON SLIDE {slide_idx} ===")
        for shape, size in picture_shapes_with_size:
            print(f"  {shape.name}: size={size:,}, >= {MIN_IMAGE_SIZE:,}? {size >= MIN_IMAGE_SIZE}")
        
        # Filter: only images larger than minimum size threshold
        large_only = [(shape, size) for shape, size in picture_shapes_with_size if size >= MIN_IMAGE_SIZE]
        
        print(f"\n=== LARGE IMAGES (>{MIN_IMAGE_SIZE:,}) ===")
        for shape, size in large_only:
            print(f"  {shape.name}: {size:,}")
        
        # Sort by image number in name (Image 0, Image 1, etc.)
        def get_image_number(shape):
            """Extract number from image name like 'Image 0', 'Image 2', etc."""
            import re
            match = re.search(r'(\d+)', shape.name)
            return int(match.group(1)) if match else 999
        
        large_only.sort(key=lambda x: get_image_number(x[0]))
        
        print(f"\n=== AFTER SORTING BY NUMBER ===")
        for shape, size in large_only[:4]:
            img_num = get_image_number(shape)
            print(f"  {shape.name} (num={img_num}): {size:,}")
        
        large_picture_shapes = [shape for shape, size in large_only[:4]]
        print(f"\nWill replace {len(large_picture_shapes)} images")
    else:
        large_picture_shapes = []
    
    # Separate logos from case study images
    logo_shapes = []
    case_study_image_shapes = []
    
    for shape in picture_shapes:
        shape_name = getattr(shape, 'name', '').lower()
        if 'logo' in shape_name:
            logo_shapes.append(shape)
        elif shape.width * shape.height >= MIN_IMAGE_SIZE:
            case_study_image_shapes.append(shape)
    
    # Match picture shapes to our case study images in order
    image_placeholders_ordered = [
        ('case_study_1_image', None),
        ('case_study_2_image', None),
        ('case_study_3_image', None),
        ('case_study_4_image', None),
    ]
    
    # Match logo placeholders
    logo_placeholders_ordered = [
        'logo1', 'logo2', 'logo3', 'logo4'
    ]
    
    # Find which placeholders have image values
    active_images = []
    for key, _ in image_placeholders_ordered:
        if key in placeholders and isinstance(placeholders[key], str) and placeholders[key].startswith("images/"):
            active_images.append((key, placeholders[key]))
    
    # Get standard image size and alignment from placeholders
    # Use FIXED DIMENSIONS only for slide 0 (overview with 4 case studies)
    # For other slides (1, 2), keep original template dimensions
    if slide_idx == 0:
        standard_width = STANDARD_IMAGE_WIDTH_EMU  # Fixed: 5.86 cm
        standard_height = STANDARD_IMAGE_HEIGHT_EMU  # Fixed: 6.0 cm
    else:
        # For slides 1 and 2, use original template dimensions
        standard_width = large_picture_shapes[0].width if large_picture_shapes else STANDARD_IMAGE_WIDTH_EMU
        standard_height = large_picture_shapes[0].height if large_picture_shapes else STANDARD_IMAGE_HEIGHT_EMU
    
    standard_top = None  # For vertical alignment
    standard_lefts = []  # For horizontal positions
    
    if large_picture_shapes:
        
        # Find minimum top position for vertical alignment
        standard_top = min(shape.top for shape in large_picture_shapes)
        
        # Collect all left positions (preserve horizontal spacing)
        standard_lefts = [shape.left for shape in large_picture_shapes]
        
        size_note = "(FIXED 6cm x 5.86cm)" if slide_idx == 0 else "(template size)"
        print(f"\n=== STANDARD IMAGE SIZE: {standard_width/914400:.2f}in x {standard_height/914400:.2f}in {size_note} ===")
        print(f"=== VERTICAL ALIGNMENT: Top={standard_top/914400:.2f}in (minimum) ===")
        lefts_str = ", ".join([f"{left/914400:.2f}" for left in standard_lefts])
        print(f"=== HORIZONTAL POSITIONS: [{lefts_str}] ===")
    
    # Match the first N large picture shapes to our N active images
    print(f"\n=== MATCHING {len(active_images)} IMAGES TO {len(large_picture_shapes)} SHAPES ===")
    for idx, (key, image_value) in enumerate(active_images):
        print(f"  Case {idx+1} ({key}): {image_value}")
        if idx < len(large_picture_shapes):
            shape = large_picture_shapes[idx]
            image_path = project_root / image_value
            print(f"    >> Will replace '{shape.name}'")
            if image_path.exists():
                # Use standard dimensions and alignment for ALL images
                # Preserve horizontal position, but align vertically
                placeholder_left = standard_lefts[idx] if idx < len(standard_lefts) else shape.left
                placeholder_top = standard_top  # All images aligned to same top
                
                # All images use the SAME standard size and vertical position
                images_to_add.append({
                    'path': str(image_path),
                    'left': placeholder_left,
                    'top': placeholder_top,
                    'width': standard_width,
                    'height': standard_height
                })
                
                print(f"    Position: Left={placeholder_left/914400:.2f}in, Top={placeholder_top/914400:.2f}in")
                size_note = "FIXED 6cm x 5.86cm" if slide_idx == 0 else "template size"
                print(f"    Size: {standard_width/914400:.2f}in x {standard_height/914400:.2f}in ({size_note})")
                
                # Mark shape for removal
                shapes_to_remove.append(shape)
    
    # Process text replacements for all shapes
    for shape in slide.shapes:
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                _process_text_shape(sub_shape, placeholders, slide_idx)
        else:
            _process_text_shape(shape, placeholders, slide_idx)
    
    # Remove shapes that had image placeholders FIRST
    for shape in shapes_to_remove:
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except Exception as e:
            print(f"Warning: Could not remove shape: {e}")
    
    # Add images AFTER removing old ones
    # For slide 0: crop images to match aspect ratio before inserting
    # For other slides: insert images directly (no cropping needed)
    temp_files = []  # Track temporary files for cleanup
    for image_info in images_to_add:
        try:
            image_path = image_info['path']
            
            # Crop images on slide 0 to match 5.86cm x 6.0cm aspect ratio
            if slide_idx == 0:
                cropped_path = _crop_image_to_aspect_ratio(image_path, 5.86, 6.0)
                temp_files.append(cropped_path)
                image_path = cropped_path
            
            pic = slide.shapes.add_picture(
                image_path,
                image_info['left'],
                image_info['top'],
                image_info['width'],
                image_info['height']
            )
            # Send image to back so it doesn't cover text
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)  # Insert near beginning (after background)
        except Exception as e:
            print(f"Warning: Could not insert image {image_info['path']}: {e}")
    
    # Clean up temporary cropped images
    for temp_file in temp_files:
        try:
            os.unlink(temp_file)
        except:
            pass
    
    # Handle logo replacements
    if logo_shapes:
        print(f"\n=== LOGO REPLACEMENT ON SLIDE {slide_idx} ===")
        print(f"Found {len(logo_shapes)} logo placeholders")
        
        # Sort logos by name (logo1, logo2, etc.)
        logo_shapes.sort(key=lambda s: getattr(s, 'name', ''))
        
        for idx, logo_placeholder in enumerate(logo_placeholders_ordered):
            if logo_placeholder in placeholders and placeholders[logo_placeholder]:
                logo_path_str = placeholders[logo_placeholder]
                logo_full_path = project_root / logo_path_str
                
                if logo_full_path.exists() and idx < len(logo_shapes):
                    old_logo = logo_shapes[idx]
                    logo_left = old_logo.left
                    logo_top = old_logo.top
                    logo_width = old_logo.width
                    logo_height = old_logo.height
                    
                    print(f"  Replacing {old_logo.name} with {logo_path_str}")
                    print(f"    Position: left={logo_left/914400:.2f}in, top={logo_top/914400:.2f}in")
                    print(f"    Size: {logo_width/914400:.2f}in x {logo_height/914400:.2f}in")
                    
                    # Remove old logo
                    try:
                        sp = old_logo.element
                        sp.getparent().remove(sp)
                    except Exception as e:
                        print(f"    Warning: Could not remove logo: {e}")
                    
                    # Convert SVG to PNG before inserting
                    try:
                        # Calculate pixel dimensions from EMUs (assuming 96 DPI)
                        width_inches = logo_width / 914400
                        height_inches = logo_height / 914400
                        width_px = int(width_inches * 96 * 3)  # 3x for better quality
                        height_px = int(height_inches * 96 * 3)
                        
                        png_path = _convert_svg_to_png(str(logo_full_path), width_px, height_px)
                        if png_path:
                            temp_files.append(png_path)  # Add to cleanup list
                            new_logo = slide.shapes.add_picture(
                                png_path,
                                logo_left,
                                logo_top,
                                logo_width,
                                logo_height
                            )
                            print(f"    Successfully inserted logo (SVG->PNG)")
                        else:
                            print(f"    Warning: Could not convert SVG to PNG")
                    except Exception as e:
                        print(f"    Warning: Could not insert logo {logo_path_str}: {e}")
    
    # Resize grey boxes to hug category text (after all text replacements are done)
    _resize_grey_boxes(slide, slide_idx)
    
    # Align metrics with their labels (tight spacing)
    _align_metrics_with_labels(slide, slide_idx)
    
    # Align titles, lines, and descriptions vertically (ONLY on slide 0)
    if slide_idx == 0:
        _align_title_line_description(slide, slide_idx)


def _process_text_shape(shape, placeholders: Dict[str, str], slide_idx: int = 0):
    """Process text replacements for a shape."""
    if hasattr(shape, "text_frame"):
        _replace_in_text_frame(shape.text_frame, placeholders, slide_idx)
    
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_in_text_frame(cell.text_frame, placeholders, slide_idx)


def _replace_in_text_frame(text_frame, placeholders: Dict[str, str], slide_idx: int = 0):
    """Replace placeholders in text frame."""
    for paragraph in text_frame.paragraphs:
        # Get full paragraph text
        full_text = paragraph.text
        
        # Check if any placeholder exists in paragraph
        has_placeholder = False
        is_infrastructure = False
        is_metric_label_on_slide2 = False
        is_metric_number_on_slide2 = False
        
        for key in placeholders.keys():
            pattern = f"{{{{{key}}}}}"
            if pattern in full_text:
                value = placeholders[key]
                # Skip image placeholders (they're handled separately)
                if isinstance(value, str) and value.startswith("images/"):
                    continue
                has_placeholder = True
                
                # Check if this is an infrastructure category
                if "category" in key and value == "INFRASTRUCTURE":
                    is_infrastructure = True
                
                # Check if this is a metric label on slide 2
                if slide_idx == 1 and "metric_label_case_study" in key:
                    is_metric_label_on_slide2 = True
                
                # Check if this is a metric number (n1, n2, n3, n4) on slide 2
                if slide_idx == 1 and key in ['n1', 'n2', 'n3', 'n4']:
                    is_metric_number_on_slide2 = True
                
                full_text = full_text.replace(pattern, str(value))
        
        # If we found placeholders, replace the entire paragraph text
        if has_placeholder and paragraph.runs:
            from pptx.util import Pt
            
            # Get current font size BEFORE any changes
            current_size = paragraph.runs[0].font.size
            if current_size:
                current_pt = current_size.pt
            else:
                current_pt = 18  # Default size if not set
            
            # Clear all runs except first
            for run in paragraph.runs[1:]:
                run.text = ""
            
            # Set first run to full replaced text
            paragraph.runs[0].text = full_text
            
            # AFTER setting text, apply font size (this is critical!)
            if is_infrastructure:
                # Infrastructure category: reduce by 1
                new_size = current_pt - 1
                paragraph.runs[0].font.size = Pt(new_size)
            elif is_metric_number_on_slide2:
                # Metric number (n1, n2, n3, n4) on slide 2: set to 30pt
                paragraph.runs[0].font.size = Pt(30)
            elif is_metric_label_on_slide2:
                # Metric label on slide 2: set to exactly 11pt
                paragraph.runs[0].font.size = Pt(11)
            else:
                # Preserve original size
                if current_size:
                    paragraph.runs[0].font.size = current_size


def add_company_context(placeholders: Dict[str, str], company_name: str, company_description: str):
    """Add company metadata to placeholders using template config."""
    placeholders.update({
        TEMPLATE_CONFIG["company_name"]: company_name,
        TEMPLATE_CONFIG["company_description"]: company_description,
        TEMPLATE_CONFIG["generation_date"]: datetime.now().strftime("%Y-%m-%d"),
        TEMPLATE_CONFIG["slide_title"]: f"Selected Case Studies — {company_name}"
    })


def _wrap_text_to_lines(text: str, width_emu: int, font_size_pt: float, max_lines: int = MAX_TITLE_LINES,
                        char_width_factor: float = TITLE_CHAR_WIDTH_FACTOR) -> tuple[list[str], bool]:
    """Wrap text to a maximum width using an approximate character width model."""
    if not text:
        return [""], False

    width_pt = width_emu / EMU_PER_POINT
    font_size_pt = font_size_pt or 11
    char_limit = max(6, int(width_pt / (font_size_pt * char_width_factor)))

    wrapped_lines: list[str] = []
    for segment in text.split('\n'):
        segment = segment.strip()
        if not segment:
            wrapped_lines.append("")
            continue
        lines = textwrap.wrap(segment, width=char_limit, break_long_words=False)
        wrapped_lines.extend(lines if lines else [segment])

    if not wrapped_lines:
        wrapped_lines = [""]

    truncated = len(wrapped_lines) > max_lines
    if truncated:
        wrapped_lines = wrapped_lines[:max_lines]
        wrapped_lines[-1] = wrapped_lines[-1].rstrip(" .,;") + "…"

    return wrapped_lines, truncated


def _calculate_title_height(line_count: int) -> int:
    """Return title height in EMUs based on the number of lines."""
    if line_count <= 0:
        return TITLE_HEIGHT_BASE_EMU
    return TITLE_HEIGHT_BASE_EMU + max(0, line_count - 1) * TITLE_HEIGHT_DELTA_EMU


def _calculate_line_width(image_width: int) -> int:
    """Calculate line width as a fixed percentage of image width (matching example1.pptx).
    
    Args:
        image_width: Width of the image in EMUs
    
    Returns:
        Width for the line shape in EMUs
    """
    return int(image_width * LINE_WIDTH_PERCENT)


def _crop_image_to_aspect_ratio(image_path: str, target_width_cm: float, target_height_cm: float) -> str:
    """Crop image to match target aspect ratio using center crop.
    
    Args:
        image_path: Path to the original image
        target_width_cm: Target width in cm (5.86)
        target_height_cm: Target height in cm (6.0)
    
    Returns:
        Path to the cropped temporary image file
    """
    target_aspect = target_width_cm / target_height_cm  # 5.86 / 6.0 = 0.9767
    
    # Open image
    img = Image.open(image_path)
    img_width, img_height = img.size
    img_aspect = img_width / img_height
    
    # Calculate crop dimensions to match target aspect ratio
    if img_aspect > target_aspect:
        # Image is too wide, crop horizontally
        new_width = int(img_height * target_aspect)
        new_height = img_height
        left = (img_width - new_width) // 2
        top = 0
        right = left + new_width
        bottom = img_height
    else:
        # Image is too tall, crop vertically
        new_width = img_width
        new_height = int(img_width / target_aspect)
        left = 0
        top = (img_height - new_height) // 2
        right = img_width
        bottom = top + new_height
    
    # Crop image
    cropped_img = img.crop((left, top, right, bottom))
    
    # Save to temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.jpg')
    os.close(temp_fd)
    cropped_img.save(temp_path, 'JPEG', quality=95)
    
    return temp_path


def _convert_svg_to_png(svg_path: str, width_px: int = 200, height_px: int = 200) -> str:
    """Convert SVG to PNG using svglib + reportlab (works on all platforms).
    
    Args:
        svg_path: Path to SVG file
        width_px: Target width in pixels
        height_px: Target height in pixels
    
    Returns:
        Path to temporary PNG file or None if conversion fails
    """
    try:
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM
        
        # Convert SVG to ReportLab drawing
        drawing = svg2rlg(svg_path)
        
        if not drawing:
            print(f"Warning: Could not parse SVG {svg_path}")
            return None
        
        # Scale drawing to desired size
        scale_x = width_px / drawing.width if drawing.width > 0 else 1
        scale_y = height_px / drawing.height if drawing.height > 0 else 1
        scale = min(scale_x, scale_y)  # Maintain aspect ratio
        
        drawing.width = drawing.width * scale
        drawing.height = drawing.height * scale
        drawing.scale(scale, scale)
        
        # Create temporary PNG file
        temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
        os.close(temp_fd)
        
        # Render to PNG at 300 DPI for high quality
        renderPM.drawToFile(drawing, temp_path, fmt='PNG', dpi=300)
        return temp_path
    except Exception as e:
        print(f"Warning: Could not convert SVG {svg_path}: {e}")
        return None
